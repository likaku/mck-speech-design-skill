#!/usr/bin/env python3
"""Validate Musk-style PPTX output end-to-end."""
import zipfile, os, json
from pptx import Presentation

print('=== 文件检查 ===')
files = {
    'PPTX (musk)': 'DeepSeek_V4_Weekly_Report_musk.pptx',
    'JSON (musk)': 'notes_musk.json',
    'Markdown':    'DeepSeek_V4_Speech_Musk.md',
    'Word (musk)': 'DeepSeek_V4_Speech_Musk.docx',
}
for label, path in files.items():
    if os.path.exists(path):
        size = os.path.getsize(path)
        print(f'  {label}: {size/1024:.1f} KB ✅')
    else:
        print(f'  {label}: ❌ 不存在')

print('\n=== PPTX 结构验证 ===')
pptx_path = 'DeepSeek_V4_Weekly_Report_musk.pptx'
with zipfile.ZipFile(pptx_path, 'r') as z:
    bad = z.testzip()
    names = z.namelist()
    notes_xmls = [n for n in names if 'notesSlide' in n and n.endswith('.xml')]
    notes_rels = [n for n in names if 'notesSlide' in n and n.endswith('.rels')]
    print(f'  ZIP 完整: {bad is None} ✅')
    print(f'  文件数: {len(names)}')
    print(f'  notes XMLs: {len(notes_xmls)}')
    print(f'  notes rels: {len(notes_rels)}')

print('\n=== 备注内容验证 ===')
prs = Presentation(pptx_path)
with open('notes_musk.json', 'r', encoding='utf-8') as f:
    expected = json.load(f)['slide_notes']

ok = fail = 0
for i, slide in enumerate(prs.slides, 1):
    key = str(i)
    if slide.has_notes_slide:
        text = slide.notes_slide.notes_text_frame.text
        exp_len = len(expected.get(key, ''))
        match = len(text) > 0 and abs(len(text) - exp_len) < 5
        status = '✅' if match else '⚠️'
        if match: ok += 1
        else: fail += 1
        musk_markers = ['Basically', 'insane', 'first principles', 'What people']
        found = [m for m in musk_markers if m.lower() in text.lower()]
        print(f'  Slide {i:2d}: {len(text):>4d} chars {status} | Musk词: {len(found)}/4 {found}')
    else:
        fail += 1
        print(f'  Slide {i:2d}: ❌ 无备注')

print()
if fail == 0:
    print(f'✅ 全部 {ok} 页验证通过！')
else:
    print(f'⚠️ {fail} 页有问题')

print('\n=== 马斯克风格密度统计 ===')
all_text = ''.join(expected.values())
markers = {
    'Basically': all_text.lower().count('basically'),
    'insane': all_text.lower().count('insane'),
    'first principles': all_text.lower().count('first principles'),
    "What people don't realize": all_text.lower().count("what people don"),
    'Um/um': all_text.count('um...') + all_text.count('Um...'),
    'epic': all_text.lower().count('epic'),
    'blows my mind': all_text.lower().count('blows my mind'),
    'The future': all_text.lower().count('the future'),
}
print(f'  总字符数: {len(all_text)}')
for k, v in markers.items():
    print(f'  "{k}": {v} 次')
