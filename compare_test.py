#!/usr/bin/env python3
"""Compare v1 (old) and v2 (new) PPTX internal XML structure."""
import zipfile

files = {
    'v1(旧)': 'DeepSeek_V4_Weekly_Report_with_notes.pptx',
    'v2(新)': 'DeepSeek_V4_Weekly_Report_with_notes_v2.pptx',
}

check_paths = [
    'ppt/slides/_rels/slide1.xml.rels',
    'ppt/notesSlides/notesSlide1.xml',
    '[Content_Types].xml',
]

for label, pptx_path in files.items():
    print(f'\n{"#"*80}')
    print(f'# {label}: {pptx_path}')
    print(f'{"#"*80}')
    with zipfile.ZipFile(pptx_path, 'r') as z:
        for cp in check_paths:
            if cp in z.namelist():
                content = z.read(cp).decode('utf-8')
                print(f'\n--- {cp} (前500字符) ---')
                print(content[:500])
            else:
                print(f'\n--- {cp}: NOT FOUND ---')

# 额外检查：用 python-pptx 读取两个文件的备注
print(f'\n{"#"*80}')
print('# python-pptx 备注读取验证')
print(f'{"#"*80}')

from pptx import Presentation

for label, pptx_path in files.items():
    print(f'\n--- {label} ---')
    try:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides, 1):
            if slide.has_notes_slide:
                text = slide.notes_slide.notes_text_frame.text
                preview = text[:60].replace('\n', ' ')
                print(f'  Slide {i:2d}: ✅ {len(text)} chars — "{preview}..."')
            else:
                print(f'  Slide {i:2d}: ❌ no notes')
    except Exception as e:
        print(f'  ❌ Error: {e}')

# Word 文档验证
print(f'\n{"#"*80}')
print('# Word 文档验证')
print(f'{"#"*80}')

from docx import Document
import os

for docx_file in ['DeepSeek_V4_Speech.docx', 'DeepSeek_V4_Speech_v2.docx']:
    if not os.path.exists(docx_file):
        print(f'\n{docx_file}: ❌ 文件不存在')
        continue
    try:
        doc = Document(docx_file)
        para_count = len(doc.paragraphs)
        table_count = len(doc.tables)
        # 检查字体
        fonts = set()
        for p in doc.paragraphs[:10]:
            for r in p.runs:
                if r.font.name:
                    fonts.add(r.font.name)
        size_kb = os.path.getsize(docx_file) / 1024
        print(f'\n{docx_file}: ✅ {size_kb:.1f}KB | {para_count} paragraphs | {table_count} tables | fonts: {fonts}')
    except Exception as e:
        print(f'\n{docx_file}: ❌ Error: {e}')
